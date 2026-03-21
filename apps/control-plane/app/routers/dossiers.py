"""Dossier router — get dossier, PDF, PPTX, markdown, ZIP, search inputs, citations."""
from __future__ import annotations

import io
import logging
import zipfile
from uuid import uuid4

from fastapi import APIRouter, Depends, Response
from sqlalchemy.orm import Session

from ..auth import AuthenticatedUser, require_user
from ..database import get_session
from ..dependencies import file_search_client
from ..errors import AppError
from ..models import (
    CitationItem,
    DossierReadModel,
    SearchMissionInputsRequest,
    SearchMissionInputsResponse,
)
from ..repository import ControlPlaneRepository
from ..services.render_service import md_to_pdf_bytes

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/missions", tags=["dossiers"])

DOC_ID_TO_ZIP_PATH = {
    "implementation_plan": "CLAUDE.md",
    "user_guide": "user_guide.md",
    "executive_summary": "executive_summary.md",
    "vision_produit": "01-strategy/vision_produit.md",
    "problem_statement": "01-strategy/problem_statement.md",
    "icp_personas": "01-strategy/icp_personas.md",
    "value_proposition": "01-strategy/value_proposition.md",
    "business_model": "02-business/business_model.md",
    "pricing_strategy": "02-business/pricing_strategy.md",
    "market_analysis": "02-business/market_analysis.md",
    "scope_document": "03-product/scope_document.md",
    "mvp_definition": "03-product/mvp_definition.md",
    "prd": "03-product/prd.md",
    "user_stories": "03-product/user_stories.md",
    "feature_specs": "03-product/feature_specs.md",
    "architecture": "04-technical/architecture.md",
    "tech_stack": "04-technical/tech_stack.md",
    "data_model": "04-technical/data_model.md",
    "api_spec": "04-technical/api_spec.md",
    "nfr_security": "04-technical/nfr_security.md",
    "ux_principles": "05-design/ux_principles.md",
    "information_architecture": "05-design/information_architecture.md",
    "design_system": "05-design/design_system.md",
    "dossier_consolide": "06-synthesis/dossier_consolide.md",
}


@router.get("/{mission_id}/dossier", response_model=DossierReadModel)
async def get_dossier(
    mission_id: str,
    user: AuthenticatedUser = Depends(require_user),
    session: Session = Depends(get_session),
):
    repository = ControlPlaneRepository(session)
    dossier = repository.get_dossier_for_user(user.id, mission_id)
    if not dossier:
        raise AppError.not_found("dossier_not_found", "Dossier not found.")
    return dossier


@router.get("/{mission_id}/dossier/pdf")
async def get_dossier_pdf(
    mission_id: str,
    user: AuthenticatedUser = Depends(require_user),
    session: Session = Depends(get_session),
):
    """Download all docs as a ZIP of individual PDF files."""
    repository = ControlPlaneRepository(session)
    dossier = repository.get_dossier_for_user(user.id, mission_id)
    if not dossier:
        raise AppError.not_found("dossier_not_found", "Dossier not found.")

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for section in dossier.sections:
            zip_path = DOC_ID_TO_ZIP_PATH.get(section.id)
            if zip_path:
                pdf_path = zip_path.replace(".md", ".pdf")
            else:
                pdf_path = f"{section.id}.pdf"
            try:
                pdf_bytes = md_to_pdf_bytes(section.title, section.content)
                zf.writestr(pdf_path, pdf_bytes)
            except Exception as exc:
                logger.warning("PDF generation failed for %s: %s", section.id, exc)
                zf.writestr(pdf_path.replace(".pdf", ".md"), f"# {section.title}\n\n{section.content}")
    buffer.seek(0)

    repository.create_export(
        export_id=str(uuid4()),
        mission_id=mission_id,
        format="PDF",
    )

    return Response(
        content=buffer.getvalue(),
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="cadris-{mission_id}-pdf.zip"'},
    )


@router.get("/{mission_id}/dossier/pptx")
async def get_dossier_pptx(
    mission_id: str,
    user: AuthenticatedUser = Depends(require_user),
    session: Session = Depends(get_session),
):
    """Generate a PowerPoint presentation from the dossier sections."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    repository = ControlPlaneRepository(session)
    dossier = repository.get_dossier_for_user(user.id, mission_id)
    if not dossier:
        raise AppError.not_found("dossier_not_found", "Dossier not found.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = dossier.title
    if slide.placeholders[1]:
        slide.placeholders[1].text = dossier.summary

    for section in dossier.sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section.title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        content = section.content[:2000] if len(section.content) > 2000 else section.content
        tf.text = content
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(14)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Synthese"
    body = slide.placeholders[1]
    body.text_frame.text = dossier.summary

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    repository.create_export(
        export_id=str(uuid4()),
        mission_id=mission_id,
        format="PPTX",
    )

    return Response(
        content=buffer.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="cadris-{mission_id}.pptx"'},
    )


@router.get("/{mission_id}/dossier/markdown")
async def get_dossier_markdown(
    mission_id: str,
    user: AuthenticatedUser = Depends(require_user),
    session: Session = Depends(get_session),
):
    """Download all docs as a ZIP of .md files, organized by folder."""
    repository = ControlPlaneRepository(session)
    dossier = repository.get_dossier_for_user(user.id, mission_id)
    if not dossier:
        raise AppError.not_found("dossier_not_found", "Dossier not found.")

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for section in dossier.sections:
            zip_path = DOC_ID_TO_ZIP_PATH.get(section.id)
            raw = section.content.strip()
            if raw.startswith("#"):
                content = raw
            else:
                content = f"# {section.title}\n\n{raw}"
            if zip_path:
                zf.writestr(zip_path, content)
            else:
                zf.writestr(f"{section.id}.md", content)
    buffer.seek(0)

    repository.create_export(
        export_id=str(uuid4()),
        mission_id=mission_id,
        format="Markdown",
    )

    return Response(
        content=buffer.getvalue(),
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="cadris-{mission_id}-md.zip"'},
    )


@router.get("/{mission_id}/dossier/zip")
async def get_dossier_zip(
    mission_id: str,
    user: AuthenticatedUser = Depends(require_user),
    session: Session = Depends(get_session),
):
    repository = ControlPlaneRepository(session)
    dossier = repository.get_dossier_for_user(user.id, mission_id)
    if not dossier:
        raise AppError.not_found("dossier_not_found", "Dossier not found.")

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for section in dossier.sections:
            zip_path = DOC_ID_TO_ZIP_PATH.get(section.id)
            if zip_path:
                raw = section.content.strip()
                content = raw if raw.startswith("#") else f"# {section.title}\n\n{raw}"
                zf.writestr(zip_path, content)

    buffer.seek(0)

    repository.create_export(
        export_id=str(uuid4()),
        mission_id=mission_id,
        format="Zip",
    )

    return Response(
        content=buffer.getvalue(),
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="cadris-{mission_id}.zip"'},
    )


@router.post(
    "/{mission_id}/inputs/search",
    response_model=SearchMissionInputsResponse,
)
async def search_mission_inputs(
    mission_id: str,
    payload: SearchMissionInputsRequest,
    user: AuthenticatedUser = Depends(require_user),
    session: Session = Depends(get_session),
):
    repository = ControlPlaneRepository(session)
    mission = repository.get_mission_for_user(user.id, mission_id)
    if not mission:
        raise AppError.not_found("mission_not_found", "Mission not found.")

    if file_search_client is None:
        return SearchMissionInputsResponse(results=[])

    vector_store_id = repository.get_vector_store_id_for_mission(mission_id)
    if not vector_store_id:
        return SearchMissionInputsResponse(results=[])

    file_id_map: dict[str, str] = {}
    for inp in repository.list_mission_inputs(mission_id):
        if inp.openai_file_id:
            file_id_map[inp.openai_file_id] = inp.id

    search_results = file_search_client.search(
        vector_store_id=vector_store_id,
        query=payload.query,
        max_results=payload.max_results,
        file_id_to_input_id=file_id_map,
    )

    citations: list[CitationItem] = []
    for result in search_results:
        citation = repository.create_citation(
            citation_id=f"cite_{uuid4().hex[:10]}",
            mission_id=mission_id,
            input_id=result.input_id,
            agent_code="search",
            excerpt=result.excerpt,
            locator=result.locator,
            score=result.score,
        )
        citation.display_name = result.filename
        citations.append(citation)

    return SearchMissionInputsResponse(results=citations)


@router.get("/{mission_id}/citations", response_model=list[CitationItem])
async def list_mission_citations(
    mission_id: str,
    user: AuthenticatedUser = Depends(require_user),
    session: Session = Depends(get_session),
):
    repository = ControlPlaneRepository(session)
    mission = repository.get_mission_for_user(user.id, mission_id)
    if not mission:
        raise AppError.not_found("mission_not_found", "Mission not found.")
    return repository.list_citations_for_mission(mission_id)
